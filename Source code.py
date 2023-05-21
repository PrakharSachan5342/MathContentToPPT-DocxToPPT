!pip install python-docx #download this pacakages before running the code
!pip install python-pptx



import docx
from pptx import Presentation
from google.colab import files

# Load the .docx file
docx_file = docx.Document("mathcontent 1.docx")  #docx file name here like i have given mathcontent 1

# Create a new PowerPoint presentation
ppt = Presentation()

# Iterate through the paragraphs in the .docx file
for paragraph in docx_file.paragraphs:
    # Check if the paragraph contains mathematical content
    if any(keyword in paragraph.text.lower() for keyword in ["equation", "graph", "diagram"]):
        # Extract the content from the paragraph
        content = paragraph.text

        # Add logic to process equations
        if "equation" in paragraph.text.lower():
            equation = process_equation(content)  # Function to parse and process the equation
            equation_result = calculate_result(equation)  # Function to calculate result

            # Add the equation and result to the slide
            slide_layout = ppt.slide_layouts[1]
            slide = ppt.slides.add_slide(slide_layout)

            # Use fixed values for positioning and sizing
            left = 100
            top = 200
            width = 400
            height = 300

            content_box = slide.shapes.add_textbox(left, top, width, height)
            content_frame = content_box.text_frame
            p = content_frame.add_paragraph()
            p.text = f"Equation: {equation}\nResult: {equation_result}"

        # Add logic to handle graphs or diagrams
        if "graph" in paragraph.text.lower() or "diagram" in paragraph.text.lower():
            graph_data = extract_graph_data(content)  # Function to extract graph data
            graph_image = generate_graph_image(graph_data)  # Function to generate graph image

            # Add the graph image to the slide
            slide_layout = ppt.slide_layouts[1]
            slide = ppt.slides.add_slide(slide_layout)

            # Use calculated values for positioning and sizing
            slide_width = ppt.slide_width
            slide_height = ppt.slide_height

            width = slide_width * 0.8
            height = slide_height * 0.6

            left = (slide_width - width) / 2
            top = (slide_height - height) / 2

            content_box = slide.shapes.add_picture(graph_image, left, top, width, height)

# Save the PowerPoint presentation
ppt.save("output.pptx")

# Download the output file
files.download("output.pptx")
